import os
import logging
import json
import time
import re
from datetime import datetime
from config import Config
from sqlalchemy.exc import OperationalError, DBAPIError

logger = logging.getLogger(__name__)

# Phase B Stage B8 (2026-05-25). Hierarchy levels for structural headers.
# Used by extract_section_headers + get_section_path_at_position to build a
# multi-level section_path for each chunk (e.g. ["Section II", "Part b"]
# instead of just "Section II:"). Lower-numbered levels are OUTER:
#
#   level 0: positional markers (Page N, Slide N) — never popped by other headers
#   level 1: Section (Roman or numeric)
#   level 2: Part (A-Z)
#   level 3: Problem / Question / Exercise — innermost
#
# When a new header at level L is opened, all entries at level >= L are popped
# from the stack (we're starting a new section AT that level). Level-0 markers
# are an exception: they only pop other level-0 markers, since a new Page or
# Slide doesn't reset the doc's outline.
_HEADER_PATTERNS = [
    # (regex_with_one_capture_group, level, display_normalizer)
    (r'(?:^|\n)(---\s*Page\s+\d+\s*---)', 0, lambda m: m.strip(' -').strip()),
    (r'(?:^|\n)(Slide\s+\d+)[:\s]?[^\n]{0,60}', 0, lambda m: m.strip()),
    (r'(?:^|\n)(Section\s+(?:\d+|[IVX]+))[:\s\-][^\n]{0,60}', 1, lambda m: m.strip()),
    (r'(?:^|\n)(Part\s+[A-Z])[:\s][^\n]{0,60}', 2, lambda m: m.strip()),
    (r'(?:^|\n)(Problem\s+\d+)[:\s][^\n]{0,60}', 3, lambda m: m.strip()),
    (r'(?:^|\n)(Question\s+\d+)[:\s][^\n]{0,60}', 3, lambda m: m.strip()),
    (r'(?:^|\n)(Exercise\s+\d+)[:\s][^\n]{0,60}', 3, lambda m: m.strip()),
]


def extract_section_headers(text: str) -> list:
    """
    Extract structural headers from text with their positions + hierarchy levels.

    Returns list of (header_content_start, header_text, level) tuples, sorted by
    position. The `level` is a small integer used by `get_section_path_at_position`
    to maintain a stack of open sections — see `_HEADER_PATTERNS` for the level
    assignments.

    Backwards-compat: callers that only care about (position, header_text) can
    unpack the first two elements and ignore level.

    Matches patterns like:
    - "Problem 1: Title"               → level 3
    - "Question 3:"                    → level 3
    - "Section I - Title"              → level 1
    - "Part A:"                        → level 2
    - "Slide 7: Title"                 → level 0
    - "--- Page 5 ---"                 → level 0
    """
    headers = []
    for pattern, level, normalizer in _HEADER_PATTERNS:
        for match in re.finditer(pattern, text, re.IGNORECASE):
            header_text = normalizer(match.group(1))
            header_text = re.sub(r'\s+', ' ', header_text)
            if len(header_text) > 80:
                header_text = header_text[:77] + "..."
            header_content_start = match.start(1)
            headers.append((header_content_start, header_text, level))
    headers.sort(key=lambda x: x[0])
    return headers


def get_context_for_position(headers: list, position: int) -> str:
    """Backwards-compat: returns the most recent single header text before the given position.

    Kept so B15 structural injection + the chunk_context String column keep
    working during the B8 transition. New code should call
    `get_section_path_at_position()` instead — it returns the full multi-level
    path (e.g. ["Section II", "Part b"]) which is the load-bearing structural
    signal post-B8.
    """
    context = ""
    for h in headers:
        # Accept both 2-tuple (legacy callers) and 3-tuple (B8 callers).
        header_pos, header_text = h[0], h[1]
        if header_pos <= position:
            context = header_text
        else:
            break
    return context


def get_section_path_at_position(headers: list, position: int) -> list:
    """Walk the headers up to `position` and return the current stack as a path.

    Examples:
        Section II opens → ["Section II"]
        Part b opens     → ["Section II", "Part b"]
        Question 3 opens → ["Section II", "Part b", "Question 3"]
        Part c opens     → ["Section II", "Part c"]  (pops level >= 2)
        Section III opens→ ["Section III"]            (pops level >= 1)

    Level-0 markers (Page N, Slide N) are an exception: they only pop other
    level-0 markers, not the outline (Section/Part/etc.) above them. So a Page 5
    appearing after Section II opens results in the path ["Section II", "Page 5"]
    — preserving the doc-outline context the chunk sits within.

    Returns an empty list when no headers precede the position.
    """
    stack = []  # list of (text, level)
    for h in headers:
        if len(h) < 3:
            # Legacy 2-tuple — treat as level 1.
            header_pos, header_text, level = h[0], h[1], 1
        else:
            header_pos, header_text, level = h
        if header_pos > position:
            break
        # Pop entries that are at-or-below this level. Exception: level-0
        # markers only pop other level-0 markers.
        if level == 0:
            stack = [(t, l) for (t, l) in stack if l != 0]
        else:
            stack = [(t, l) for (t, l) in stack if l < level]
        stack.append((header_text, level))
    return [t for (t, _) in stack]

def chunk_text_with_context(text: str, chunk_size: int = 800, overlap: int = 200, doc_filename: str = "") -> list:
    """
    Chunk text with boundary-aware splitting at section headers.
    Forces chunk breaks at section boundaries (Problem X, Question X, etc.) so chunks
    don't span multiple problems/sections.
    Returns list of dicts with 'text' (enriched) and 'original_text' (raw).
    """
    headers = extract_section_headers(text)

    # Cache sorted header positions once. Headers are now 3-tuples
    # (position, text, level); we only care about the position for boundary detection.
    sorted_header_positions = sorted(h[0] for h in headers)

    if len(text) <= chunk_size:
        context = get_context_for_position(headers, 0)
        section_path = get_section_path_at_position(headers, 0)
        enriched = f"[{doc_filename} > {context}] {text}" if context else f"[{doc_filename}] {text}"
        return [{"text": enriched, "original_text": text, "context": context, "section_path": section_path}]
    
    chunks = []
    start = 0
    header_idx = 0  # Track position in sorted headers
    
    while start < len(text):
        # Advance header index past any headers at or before start
        while header_idx < len(sorted_header_positions) and sorted_header_positions[header_idx] <= start:
            header_idx += 1
        next_boundary = sorted_header_positions[header_idx] if header_idx < len(sorted_header_positions) else None
        
        # Calculate tentative end position
        end = start + chunk_size
        
        # Track if we're breaking at a boundary
        breaking_at_boundary = False
        
        # If there's a section boundary within this chunk, force break there
        if next_boundary is not None and next_boundary < end and next_boundary > start:
            # End this chunk at the boundary position (not including the header)
            end = next_boundary
            breaking_at_boundary = True
        elif end < len(text):
            # No boundary in range - use natural break points as before
            break_points = [
                text.rfind('\n\n', start, end),
                text.rfind('. ', start, end),
                text.rfind('\n', start, end),
                text.rfind(' ', start, end)
            ]
            for bp in break_points:
                if bp > start + chunk_size // 2:
                    end = bp + 1
                    break
        
        # Guard against zero-length chunks or no progress
        if end <= start:
            end = min(start + chunk_size, len(text))
        
        chunk_text = text[start:end].strip()
        
        if chunk_text:
            # Get context + multi-level section_path for this chunk position
            context = get_context_for_position(headers, start)
            section_path = get_section_path_at_position(headers, start)
            if context:
                enriched = f"[{doc_filename} > {context}] {chunk_text}"
            else:
                enriched = f"[{doc_filename}] {chunk_text}"
            chunks.append({
                "text": enriched,
                "original_text": chunk_text,
                "context": context,
                "section_path": section_path,
            })
            
            # Move to next position
            if breaking_at_boundary:
                # CRITICAL: Start exactly at the boundary with NO overlap
                # This ensures the next chunk gets the correct section context
                start = next_boundary
            else:
                start = end - overlap
                if start < 0:
                    start = 0
        else:
            # Empty chunk - advance to avoid infinite loop
            if next_boundary is not None and next_boundary > start:
                start = next_boundary
            else:
                start = end if end > start else start + 1
        
        if start >= len(text):
            break
    
    return chunks

def db_commit_with_retry(db, max_retries=3, delay=1.0):
    """Commit database changes with retry logic for connection issues."""
    for attempt in range(max_retries):
        try:
            db.session.commit()
            return True
        except (OperationalError, DBAPIError) as e:
            db.session.rollback()
            if attempt < max_retries - 1:
                logger.warning(f"Database commit failed (attempt {attempt + 1}/{max_retries}): {e}")
                time.sleep(delay * (attempt + 1))
            else:
                logger.error(f"Database commit failed after {max_retries} attempts: {e}")
                raise
    return False

def sanitize_text(text: str) -> str:
    """Remove null bytes and other problematic characters that PostgreSQL cannot store."""
    if not text:
        return ""
    text = text.replace('\x00', '')
    text = ''.join(char for char in text if ord(char) >= 32 or char in '\n\r\t')
    return text

def extract_image(file_path: str) -> str:
    """Extract content from an image file using GPT-4o vision."""
    try:
        import base64
        from openai import OpenAI

        ext = file_path.rsplit('.', 1)[-1].lower()
        mime_map = {
            'jpg': 'image/jpeg', 'jpeg': 'image/jpeg',
            'png': 'image/png', 'gif': 'image/gif', 'webp': 'image/webp',
        }
        mime_type = mime_map.get(ext, 'image/jpeg')

        with open(file_path, 'rb') as f:
            img_base64 = base64.b64encode(f.read()).decode()

        client = OpenAI(api_key=Config.OPENAI_API_KEY)
        response = client.chat.completions.create(
            model=Config.VISION_MODEL,
            max_tokens=1500,
            messages=[{
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": (
                            "This is an image from an academic course. "
                            "Transcribe ALL visible content including:\n"
                            "- All text (printed or handwritten)\n"
                            "- Mathematical equations (use LaTeX notation like $x^2$)\n"
                            "- Diagrams and figures (describe in [DIAGRAM: ...] tags)\n"
                            "- Tables (preserve structure)\n"
                            "- Labels, annotations, and any other visual information\n\n"
                            "Be thorough — students will use this content to get help."
                        )
                    },
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:{mime_type};base64,{img_base64}",
                            "detail": "high"
                        }
                    }
                ]
            }]
        )
        return response.choices[0].message.content or ""
    except Exception as e:
        logger.warning(f"Image extraction failed for {file_path}: {e}")
        return ""


def extract_json(file_path: str) -> str:
    """Extract content from a JSON file as pretty-printed text."""
    try:
        import json
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            data = json.load(f)
        return json.dumps(data, indent=2, ensure_ascii=False)
    except Exception:
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            logger.warning(f"JSON extraction failed for {file_path}: {e}")
            return ""


def extract_csv(file_path: str) -> str:
    """Extract content from a CSV file as a readable table."""
    try:
        import pandas as pd
        for encoding in ('utf-8', 'latin-1', 'cp1252'):
            try:
                df = pd.read_csv(file_path, encoding=encoding)
                return df.to_string(index=False)
            except UnicodeDecodeError:
                continue
        return ""
    except Exception as e:
        logger.warning(f"CSV extraction failed for {file_path}: {e}")
        return ""


def extract_jupyter_notebook(file_path: str) -> str:
    """Extract content from a Jupyter notebook (.ipynb) file."""
    try:
        import json
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            nb = json.load(f)

        parts = []
        for cell in nb.get('cells', []):
            cell_type = cell.get('cell_type', '')
            source = ''.join(cell.get('source', []))
            if not source.strip():
                continue

            if cell_type == 'markdown':
                parts.append(source)
            elif cell_type == 'code':
                parts.append(f"```python\n{source}\n```")
                for output in cell.get('outputs', []):
                    if output.get('output_type') in ('stream', 'execute_result', 'display_data'):
                        out_text = ''.join(
                            output.get('text', []) or
                            output.get('data', {}).get('text/plain', [])
                        )
                        if out_text.strip():
                            parts.append(f"Output:\n{out_text.strip()}")

        return '\n\n'.join(parts)
    except Exception as e:
        logger.warning(f"Jupyter notebook extraction failed for {file_path}: {e}")
        return ""


def extract_text_from_file(file_path: str, heartbeat=None) -> tuple:
    """
    Extract text from file. Returns (text, page_count) - page_count is 0 for non-PDFs.

    heartbeat: optional callable the PDF path invokes periodically during
    page-by-page vision work. The indexing watchdog fails any job whose progress
    row has not moved in 5 minutes, and a long PDF can spend far longer than that
    inside a single document — the job is working, but invisibly.
    """
    ext = file_path.rsplit('.', 1)[-1].lower() if '.' in file_path else ''
    page_count = 0

    try:
        if ext == 'pdf':
            text, page_count = extract_pdf(file_path, heartbeat=heartbeat)
        elif ext == 'docx':
            text = extract_docx(file_path)
        elif ext == 'doc':
            text = extract_doc(file_path)
        elif ext in ('xlsx', 'xls'):
            text = extract_excel(file_path)
        elif ext == 'txt':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
        elif ext in ('pptx', 'ppt'):
            text = extract_pptx(file_path)
        elif ext in ('png', 'jpg', 'jpeg', 'gif', 'webp'):
            text = extract_image(file_path)
        elif ext in ('md', 'py'):
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
        elif ext == 'json':
            text = extract_json(file_path)
        elif ext == 'csv':
            text = extract_csv(file_path)
        elif ext == 'ipynb':
            text = extract_jupyter_notebook(file_path)
        else:
            logger.warning(f"Unsupported file type: {ext}")
            return "", 0

        return sanitize_text(text), page_count
    except Exception as e:
        logger.error(f"Error extracting text from {file_path}: {e}")
        return "", 0

def extract_pdf(file_path: str, heartbeat=None) -> tuple:
    """
    Extract PDF text and return (text, page_count).

    heartbeat: optional callable passed down to the page-by-page vision loops so a
    long document keeps reporting progress to the indexing watchdog.
    """
    text, page_count = _extract_pdf_pypdf2(file_path, heartbeat=heartbeat)
    if text and len(text.strip()) > 100:
        text = _supplement_pdf_with_figures(file_path, text, heartbeat=heartbeat)
        return text, page_count

    logger.info("PyPDF2 extraction insufficient, trying pdfplumber...")
    text, page_count = _extract_pdf_pdfplumber(file_path, heartbeat=heartbeat)
    if text and len(text.strip()) > 100:
        text = _supplement_pdf_with_figures(file_path, text, heartbeat=heartbeat)
        return text, page_count

    logger.info("Text extraction insufficient - attempting vision-based extraction for image/handwritten PDF...")
    text, page_count = _extract_pdf_vision(file_path, heartbeat=heartbeat)
    if text and len(text.strip()) > 50:
        return text, page_count  # vision already described figures — skip supplement

    return "", 0


def _select_figure_candidate_pages(file_path: str, heartbeat=None) -> list:
    """
    Page numbers (1-based) that plausibly contain a figure, decided from the PDF's
    own structure — no rendering and no API calls.

    A page qualifies if it has an embedded raster image or enough vector drawing
    operations to look like a plot. Tables are deliberately NOT caught: they show
    up as rects and lines rather than curves, and transcribing them is what the
    text extractor is for.

    A "page has no extractable text" condition was tried and removed. It was meant
    to catch full-page figures, but any real figure is necessarily either a raster
    or a vector drawing and is already caught — so the rule added no coverage and
    fired on genuinely blank front-matter pages. Measured: it inflated doc 161 from
    60 to 73 candidates, and 9 of the first 12 came back "No figures".

    Returns a list of page numbers on success — possibly EMPTY, which legitimately
    means "this document has no figures" and must not be confused with failure.
    Returns None if screening itself failed, so the caller can fall back to the
    pre-screening behaviour of visioning every page.
    """
    try:
        import pdfplumber

        candidates = []
        with pdfplumber.open(file_path) as pdf:
            total = len(pdf.pages)
            for i, page in enumerate(pdf.pages, 1):
                # This loop is the slowest phase of PDF extraction on a long
                # document — reading images/curves means parsing each page's
                # content stream. Measured ~125s on a 356-page textbook, longer
                # than the text extraction and the vision calls combined. Heartbeat
                # frequently or the watchdog fails the job mid-screen.
                if heartbeat and i % 25 == 0:
                    heartbeat()
                has_raster = len(page.images) > 0
                has_vector = len(page.curves) >= Config.VISION_FIGURE_CURVE_THRESHOLD
                if has_raster or has_vector:
                    candidates.append(i)

        if len(candidates) > Config.VISION_MAX_PAGES_PER_DOC:
            logger.warning(
                f"Figure supplement: {len(candidates)} candidate pages exceeds cap "
                f"{Config.VISION_MAX_PAGES_PER_DOC}; taking the first {Config.VISION_MAX_PAGES_PER_DOC}. "
                f"Figures on later pages will not be described."
            )
            candidates = candidates[:Config.VISION_MAX_PAGES_PER_DOC]

        logger.info(
            f"Figure supplement: {len(candidates)}/{total} pages look like figures "
            f"({100 * len(candidates) / total:.0f}%) for {file_path}"
        )
        return candidates

    except Exception as e:
        logger.warning(
            f"Figure-candidate screening failed ({type(e).__name__}: {e}); "
            f"falling back to every page"
        )
        return None


def _supplement_pdf_with_figures(file_path: str, text: str, heartbeat=None) -> str:
    """
    Run a figure-only vision pass over the PDF pages that plausibly contain a
    figure, splicing descriptions into the already-extracted text under the
    matching '--- Page N ---' marker.

    Pages are screened first (see _select_figure_candidate_pages) and rendered ONE
    AT A TIME. Previously this rendered the whole document up front —
    convert_from_path with no page range — so a 356-page book produced 356
    in-memory JPEGs at 200 DPI before the first API call. That cost landed before
    any vision work and was a major reason large PDFs tripped the 5-minute
    indexing watchdog.

    heartbeat: optional callable invoked every Config.VISION_HEARTBEAT_EVERY pages
    so long documents keep reporting progress. Without it the watchdog marks a job
    failed while it is still working — the job is never actually stuck.
    """
    try:
        import base64
        import re as _re
        from io import BytesIO
        from pdf2image import convert_from_path
        from openai import OpenAI

        client = OpenAI(api_key=Config.OPENAI_API_KEY)

        candidate_pages = _select_figure_candidate_pages(file_path, heartbeat=heartbeat)
        if candidate_pages is None:
            # Screening ITSELF failed — fall back to the pre-screening behaviour of
            # every page. An empty list is a different thing entirely: it means the
            # document genuinely has no figures, and we should do no vision at all.
            try:
                import pdfplumber
                with pdfplumber.open(file_path) as pdf:
                    candidate_pages = list(range(1, len(pdf.pages) + 1))
            except Exception:
                return text
        elif not candidate_pages:
            logger.info(f"Figure supplement: no figure-bearing pages in {file_path}; skipping vision")
            return text

        supplemented = text
        for seq, page_num in enumerate(candidate_pages, 1):
            try:
                # Render just this page. first_page/last_page are 1-based inclusive.
                rendered = convert_from_path(
                    file_path, dpi=200, fmt='jpeg',
                    first_page=page_num, last_page=page_num,
                    poppler_path=Config.POPPLER_PATH,
                )
                if not rendered:
                    continue
                img = rendered[0]

                if heartbeat and seq % Config.VISION_HEARTBEAT_EVERY == 0:
                    heartbeat()

                img_buffer = BytesIO()
                img.save(img_buffer, format='JPEG', quality=85)
                img_base64 = base64.b64encode(img_buffer.getvalue()).decode()

                response = client.chat.completions.create(
                    model=Config.VISION_MODEL,
                    max_tokens=1000,
                    messages=[{
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": (
                                    f"This is slide/page {page_num} of an academic lecture. "
                                    "Describe ALL visual elements: charts, graphs, figures, diagrams, maps, photos, and tables shown as images. "
                                    "For maps: list every labeled city, region, country, or location name visible. "
                                    "For tables: describe the structure and transcribe all cell contents. "
                                    "For charts/graphs: include axis labels, legend entries, data values, and trends. "
                                    "Include all text that appears inside or as part of a visual element. "
                                    "Do NOT describe regular slide titles or bullet-point body text. "
                                    "If there are no visual elements on this page, reply exactly: No figures"
                                )
                            },
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/jpeg;base64,{img_base64}",
                                    "detail": "high"
                                }
                            }
                        ]
                    }]
                )

                desc = response.choices[0].message.content or ""
                desc = desc.strip()
                if not desc or desc.lower().startswith("no figures"):
                    continue

                # Splice description after the matching "--- Page N ---" marker
                marker = f"--- Page {page_num} ---"
                if marker in supplemented:
                    supplemented = supplemented.replace(
                        marker,
                        f"{marker}\n[FIGURE: {desc}]",
                        1
                    )
                    logger.info(f"Figure supplement: added figure description for page {page_num}")
                else:
                    # Page marker missing (e.g. page had no text) — append at end
                    supplemented += f"\n\n{marker}\n[FIGURE: {desc}]"
                    logger.info(f"Figure supplement: appended figure description for page {page_num} (no marker found)")

            except Exception as page_e:
                logger.warning(f"Figure supplement failed on page {page_num}: {page_e}")
                continue

        return supplemented

    except ImportError as e:
        logger.warning(f"Figure supplement unavailable - missing dependency: {e}")
        return text
    except Exception as e:
        logger.warning(f"Figure supplement failed: {e}")
        return text


def _extract_pdf_vision(file_path: str, heartbeat=None) -> tuple:
    """
    Extract content from image-heavy/handwritten PDFs using GPT-4o vision.
    Converts each page to an image and sends to GPT-4o for transcription.
    Returns (text, page_count).

    heartbeat: optional callable invoked every Config.VISION_HEARTBEAT_EVERY pages.
    This path is already capped at 50 pages, but 50 vision calls can still exceed
    the indexing watchdog's 5-minute window on a single document.
    """
    try:
        import base64
        from io import BytesIO
        from pdf2image import convert_from_path
        from openai import OpenAI

        client = OpenAI(api_key=Config.OPENAI_API_KEY)

        images = convert_from_path(file_path, dpi=200, fmt='jpeg', poppler_path=Config.POPPLER_PATH)
        page_count = len(images)
        logger.info(f"Vision extraction: converted {page_count} pages to images")

        if page_count == 0:
            return "", 0

        max_pages = 50
        if page_count > max_pages:
            logger.warning(f"Vision extraction: PDF has {page_count} pages, limiting to {max_pages}")
            images = images[:max_pages]

        text_parts = []
        for page_num, img in enumerate(images, 1):
            try:
                if heartbeat and page_num % Config.VISION_HEARTBEAT_EVERY == 0:
                    heartbeat()

                img_buffer = BytesIO()
                img.save(img_buffer, format='JPEG', quality=85)
                img_base64 = base64.b64encode(img_buffer.getvalue()).decode()
                img_size_kb = len(img_buffer.getvalue()) / 1024
                logger.info(f"Vision extraction: processing page {page_num}/{len(images)} ({img_size_kb:.0f}KB)")

                response = client.chat.completions.create(
                    model=Config.VISION_MODEL,
                    max_tokens=1500,
                    messages=[{
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": (
                                    f"This is page {page_num} of a lecture document. "
                                    "Transcribe ALL visible content including:\n"
                                    "- Handwritten text (preserve exact wording)\n"
                                    "- Printed text\n"
                                    "- Mathematical equations and formulas (use LaTeX notation)\n"
                                    "- Diagrams and illustrations (describe them in [DIAGRAM: ...] tags)\n"
                                    "- Tables (preserve structure)\n"
                                    "- Labels, annotations, and margin notes\n\n"
                                    "For mathematical content, use LaTeX notation like $x^2$ or $$\\int f(x) dx$$.\n"
                                    "Preserve the logical flow and structure of the content. "
                                    "If text is unclear, provide your best interpretation with [unclear] markers."
                                )
                            },
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/jpeg;base64,{img_base64}",
                                    "detail": "high"
                                }
                            }
                        ]
                    }],
                )

                page_text = response.choices[0].message.content
                if page_text and page_text.strip():
                    text_parts.append(f"--- Page {page_num} ---\n{page_text.strip()}")
                    logger.info(f"Vision extraction: page {page_num} yielded {len(page_text)} chars")
                else:
                    logger.warning(f"Vision extraction: page {page_num} returned empty content")

            except Exception as page_e:
                logger.warning(f"Vision extraction failed on page {page_num}: {page_e}")
                continue

        if not text_parts:
            return "", 0

        full_text = "\n\n".join(text_parts)
        logger.info(f"Vision extraction complete: {len(full_text)} chars from {len(text_parts)}/{page_count} pages")
        return full_text, page_count

    except ImportError as e:
        logger.warning(f"Vision extraction unavailable - missing dependency: {e}")
        return "", 0
    except Exception as e:
        logger.error(f"Vision extraction failed: {e}")
        return "", 0

def _extract_pdf_pdfplumber(file_path: str, heartbeat=None) -> tuple:
    """
    Extract PDF text using pdfplumber with total time limit. Returns (text, page_count).

    heartbeat fires every 25 pages so the indexing watchdog can see progress. This
    path is already bounded at 120s by max_total_time, but that is still a large
    share of the watchdog's 5-minute window before anything else runs.
    """
    try:
        import pdfplumber

        text_parts = []
        start_time = time.time()
        max_total_time = 120

        with pdfplumber.open(file_path) as pdf:
            total_pages = len(pdf.pages)
            for page_num, page in enumerate(pdf.pages):
                if heartbeat and page_num and page_num % 25 == 0:
                    heartbeat()
                page_start = time.time()
                try:
                    if time.time() - start_time > max_total_time:
                        logger.warning(f"pdfplumber exceeded {max_total_time}s total, stopping at page {page_num}")
                        break
                    
                    text = page.extract_text()
                    if text:
                        text_parts.append(f"--- Page {page_num + 1} ---\n{text}")
                    
                    page_time = time.time() - page_start
                    if page_time > 10:
                        logger.info(f"pdfplumber page {page_num + 1}/{total_pages} took {page_time:.1f}s")
                        
                except Exception as page_e:
                    logger.warning(f"pdfplumber failed on page {page_num + 1}: {page_e}")
                    continue
                    
        return "\n\n".join(text_parts), total_pages
    except Exception as e:
        logger.warning(f"pdfplumber extraction failed: {e}")
        return "", 0

def _extract_pdf_pypdf2(file_path: str, heartbeat=None) -> tuple:
    """
    Extract PDF text using PyPDF2. Returns (text, page_count).

    heartbeat fires every 25 pages. Measured on a 356-page textbook, this loop
    alone runs ~170s — longer than the vision supplement that follows it — so
    without a heartbeat here a large PDF can burn most of the indexing watchdog's
    5-minute window before the first progress update is even possible.
    """
    try:
        from PyPDF2 import PdfReader

        reader = PdfReader(file_path)
        text_parts = []
        page_count = len(reader.pages)
        for page_num, page in enumerate(reader.pages, 1):
            if heartbeat and page_num % 25 == 0:
                heartbeat()
            text = page.extract_text()
            if text:
                text_parts.append(f"--- Page {page_num} ---\n{text}")
        return "\n\n".join(text_parts), page_count
    except Exception as e:
        logger.warning(f"PyPDF2 extraction failed: {e}")
        return "", 0

def extract_docx(file_path: str) -> str:
    """
    Extract text from DOCX files, preserving list numbering (a, b, c, 1, 2, 3),
    tables, and document structure.
    
    Uses docx2python for comprehensive extraction that preserves:
    - Numbered/bulleted list prefixes
    - Table content
    - Text boxes
    - Footnotes/endnotes
    """
    try:
        from docx2python import docx2python
        
        with docx2python(file_path) as doc:
            text_parts = []
            
            if doc.body:
                body_text = _flatten_docx2python_content(doc.body)
                if body_text.strip():
                    text_parts.append(body_text)
            
            if doc.footnotes:
                footnotes_text = _flatten_docx2python_content(doc.footnotes)
                if footnotes_text.strip():
                    text_parts.append("\n--- Footnotes ---\n" + footnotes_text)
            
            if doc.endnotes:
                endnotes_text = _flatten_docx2python_content(doc.endnotes)
                if endnotes_text.strip():
                    text_parts.append("\n--- Endnotes ---\n" + endnotes_text)
            
            result = "\n\n".join(text_parts)
            
            if result.strip():
                logger.info(f"Successfully extracted DOCX using docx2python: {len(result)} chars")
                return result
            
    except Exception as e:
        logger.warning(f"docx2python extraction failed, falling back to python-docx: {e}")
    
    try:
        from docx import Document
        
        doc = Document(file_path)
        text_parts = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                prefix = _get_paragraph_list_prefix(para)
                if prefix:
                    text_parts.append(f"{prefix} {para.text}")
                else:
                    text_parts.append(para.text)
        
        for table in doc.tables:
            table_text = _extract_table_text(table)
            if table_text.strip():
                text_parts.append(table_text)
        
        result = "\n\n".join(text_parts)
        logger.info(f"Extracted DOCX using python-docx fallback: {len(result)} chars")
        return result
        
    except Exception as e:
        logger.error(f"DOCX extraction failed completely: {e}")
        return ""


def _flatten_docx2python_content(content) -> str:
    """
    Recursively flatten docx2python nested list structure into text.
    docx2python returns deeply nested lists representing document structure.
    """
    if isinstance(content, str):
        return content.strip()
    
    if isinstance(content, list):
        parts = []
        for item in content:
            flattened = _flatten_docx2python_content(item)
            if flattened:
                parts.append(flattened)
        return "\n".join(parts)
    
    return ""


def _get_paragraph_list_prefix(para) -> str:
    """
    Extract list numbering prefix from a paragraph using python-docx XML parsing.
    Returns prefix like 'a)', '1.', 'i)' or empty string if not a list item.
    """
    try:
        if para._element.pPr is None:
            return ""
        
        numPr = para._element.pPr.numPr
        if numPr is None or numPr.numId is None:
            return ""
        
        return ""
        
    except Exception:
        return ""


def _extract_table_text(table) -> str:
    """Extract text from a Word table, preserving structure."""
    try:
        rows_text = []
        for row in table.rows:
            cells_text = []
            for cell in row.cells:
                cell_content = cell.text.strip()
                if cell_content:
                    cells_text.append(cell_content)
            if cells_text:
                rows_text.append(" | ".join(cells_text))
        
        if rows_text:
            return "Table:\n" + "\n".join(rows_text)
        return ""
    except Exception:
        return ""


def extract_doc(file_path: str) -> str:
    """
    Extract text from older .doc files (pre-2007 Word format).
    Uses antiword command-line tool or falls back to textract.
    """
    import subprocess
    
    try:
        result = subprocess.run(
            ['antiword', file_path],
            capture_output=True,
            text=True,
            timeout=30
        )
        if result.returncode == 0 and result.stdout.strip():
            logger.info(f"Extracted .doc using antiword: {len(result.stdout)} chars")
            return result.stdout
    except FileNotFoundError:
        logger.warning("antiword not installed, trying catdoc")
    except Exception as e:
        logger.warning(f"antiword extraction failed: {e}")
    
    try:
        result = subprocess.run(
            ['catdoc', file_path],
            capture_output=True,
            text=True,
            timeout=30
        )
        if result.returncode == 0 and result.stdout.strip():
            logger.info(f"Extracted .doc using catdoc: {len(result.stdout)} chars")
            return result.stdout
    except FileNotFoundError:
        logger.warning("catdoc not installed")
    except Exception as e:
        logger.warning(f"catdoc extraction failed: {e}")
    
    logger.warning(f"Could not extract .doc file: {file_path}. Install antiword or catdoc.")
    return ""

def extract_excel(file_path: str) -> str:
    import pandas as pd
    
    xl = pd.ExcelFile(file_path)
    text_parts = []
    for sheet_name in xl.sheet_names:
        df = xl.parse(sheet_name)
        text_parts.append(f"Sheet: {sheet_name}\n{df.to_string()}")
    return "\n\n".join(text_parts)

def extract_pptx(file_path: str) -> str:
    try:
        import base64
        from io import BytesIO
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from openai import OpenAI

        prs = Presentation(file_path)
        client = OpenAI(api_key=Config.OPENAI_API_KEY)
        text_parts = []
        raster_types = {'image/jpeg', 'image/png', 'image/gif', 'image/bmp', 'image/webp'}

        def _iter_shapes(shapes):
            for shape in shapes:
                yield shape
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    yield from _iter_shapes(shape.shapes)

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            figure_descriptions = []

            for shape in _iter_shapes(slide.shapes):
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img = shape.image
                        if img.content_type not in raster_types:
                            continue
                        img_base64 = base64.b64encode(img.blob).decode()
                        response = client.chat.completions.create(
                            model=Config.VISION_MODEL,
                            max_tokens=800,
                            messages=[{
                                "role": "user",
                                "content": [
                                    {
                                        "type": "text",
                                        "text": (
                                            f"Describe this visual from slide {slide_num} of an academic lecture. "
                                            "For maps: list every labeled city, region, country, or location name visible. "
                                            "For tables: describe the structure and transcribe all cell contents. "
                                            "For charts/graphs: include axis labels, legend entries, data values, and trends. "
                                            "Include all text visible inside the image (labels, annotations, callouts). "
                                            "Be thorough — do not summarize away specific names or values."
                                        )
                                    },
                                    {
                                        "type": "image_url",
                                        "image_url": {
                                            "url": f"data:{img.content_type};base64,{img_base64}",
                                            "detail": "high"
                                        }
                                    }
                                ]
                            }]
                        )
                        desc = response.choices[0].message.content
                        if desc and desc.strip():
                            figure_descriptions.append(f"[FIGURE: {desc.strip()}]")
                            logger.info(f"PPTX: described image on slide {slide_num}")
                    except Exception as img_e:
                        logger.warning(f"PPTX: failed to describe image on slide {slide_num}: {img_e}")

                if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    logger.info(f"PPTX: found CHART shape on slide {slide_num}")
                    try:
                        chart = shape.chart
                        title = chart.chart_title.text_frame.text if chart.has_title else ""
                        chart_type = str(chart.chart_type)
                        series_data = []
                        for plot in chart.plots:
                            for series in plot.series:
                                try:
                                    vals = [v for v in (series.values or []) if v is not None]
                                    series_data.append({"name": series.name or "Series", "values": vals[:20]})
                                except Exception:
                                    pass
                        prompt = (
                            f"Describe this chart from slide {slide_num} of a lecture in 2-4 sentences. "
                            f"Chart type: {chart_type}. Title: '{title}'. Series data: {series_data}. "
                            "Include the key trend, axis interpretation, and any notable data points."
                        )
                        response = client.chat.completions.create(
                            model=Config.VISION_MODEL,
                            max_tokens=500,
                            messages=[{"role": "user", "content": prompt}]
                        )
                        desc = response.choices[0].message.content
                        if desc and desc.strip():
                            figure_descriptions.append(f"[CHART: {desc.strip()}]")
                            logger.info(f"PPTX: described chart on slide {slide_num}")
                    except Exception as chart_e:
                        import traceback
                        logger.warning(f"PPTX: failed to describe chart on slide {slide_num}: {chart_e}\n{traceback.format_exc()}")

            slide_parts = []
            if slide_text:
                slide_parts.append("\n".join(slide_text))
            slide_parts.extend(figure_descriptions)
            if not slide_parts:
                slide_parts.append("[Visual content only — chart or figure description unavailable]")

            text_parts.append(f"Slide {slide_num}:\n" + "\n".join(slide_parts))

        return "\n\n".join(text_parts)
    except ImportError:
        logger.warning("python-pptx not installed, skipping PowerPoint file")
        return ""

def extract_metadata_with_llm(text: str, filename: str) -> dict:
    from openai import OpenAI
    
    client = OpenAI(api_key=Config.OPENAI_API_KEY)
    
    preview = text[:3000] if len(text) > 3000 else text
    
    prompt = f"""Analyze this course document and extract metadata. The filename is: {filename}

Document preview:
{preview}

Extract the following information as JSON:
{{
    "doc_type": "homework" | "exam" | "lecture" | "reading" | "syllabus" | "other",
    "assignment_number": "1" | "2" | "3" | null (if applicable),
    "instructional_unit_number": 1 | 2 | 3 | null (lecture/class/week number if mentioned),
    "instructional_unit_label": "lecture" | "class" | "week" | "module" | "session" | null,
    "course_code": "MGT404" | null (if visible),
    "year": "2024" | "2025" | null (if mentioned),
    "is_solutions": true | false (whether this contains solutions/answers),
    "content_title": "The actual document title as written in the content (e.g., 'Self-Study Problem Set #2', 'Final Exam 2024'). Extract from headers/title text, not filename.",
    "section_numbering_style": "arabic" | "roman" | "mixed" | null (how major sections/problems are numbered: "1, 2, 3" = arabic, "I, II, III" = roman)
}}

IMPORTANT: Classify by the PRIMARY document type, not whether it has solutions.
- "exam_solutions.pdf" or "final_exam_with_solutions.pdf" → doc_type: "exam", is_solutions: true
- "homework_answers.pdf" or "problem_set_solutions.pdf" → doc_type: "homework", is_solutions: true
- "lecture_notes.pdf" → doc_type: "lecture"

NUMBERING STYLE: Look at how the document labels its main sections/problems:
- If sections are "Section 1", "Problem 2", "Question 3" → section_numbering_style: "arabic"
- If sections are "Section I", "Problem II", "Part III" → section_numbering_style: "roman"
- If mixed or unclear → section_numbering_style: "mixed"

Return ONLY valid JSON, no other text."""

    for attempt in range(3):
        try:
            response = client.chat.completions.create(
                model=Config.LLM_MODEL,
                messages=[{"role": "user", "content": prompt}],
                max_completion_tokens=500,
                reasoning_effort=Config.LLM_REASONING_MEDIUM
            )
            
            content = response.choices[0].message.content.strip()
            if content.startswith("```"):
                content = content.split("```")[1]
                if content.startswith("json"):
                    content = content[4:]
            
            metadata = json.loads(content)
            return metadata
        except Exception as e:
            logger.warning(f"LLM metadata extraction attempt {attempt+1} failed: {e}")
            if attempt == 2:
                logger.error(f"All LLM metadata extraction attempts failed, returning defaults")
                return {
                    "doc_type": "other",
                    "assignment_number": None,
                    "instructional_unit_number": None,
                    "instructional_unit_label": None,
                    "course_code": None,
                    "year": None,
                    "is_solutions": False,
                    "content_title": None,
                    "section_numbering_style": None
                }


# Phase A Stage 2B (research 2026-05-22). Default per-TA category seed list,
# applied at TA creation. Professor edits freely via the manage_ta UI.
# Each entry is {slug, label}: the slug is the immutable internal identifier
# (used in Document.doc_category + DocumentChunk.doc_category); the label is
# the mutable display string (rename-safe — see Library Drift, arXiv 2605.19576).
DEFAULT_DOC_CATEGORIES = [
    {"slug": "lectures", "label": "Lectures"},
    {"slug": "readings", "label": "Readings"},
    {"slug": "homeworks", "label": "Homeworks"},
    {"slug": "problem_sets", "label": "Problem Sets"},
    {"slug": "quizzes", "label": "Quizzes"},
    {"slug": "labs", "label": "Labs"},
    {"slug": "exams", "label": "Exams"},
    {"slug": "syllabus", "label": "Syllabus"},
    {"slug": "reference_materials", "label": "Reference Materials"},
    {"slug": "extra_problems", "label": "Extra Problems"},
    {"slug": "solutions", "label": "Solutions"},
    {"slug": "other", "label": "Other"},
]


def normalize_category_slug(value: str) -> str:
    """Normalize a free-form category label into a stable internal slug.

    Rules (research 2026-05-22, Q5 failure-mode mitigation):
    - lowercase
    - internal whitespace → underscores
    - strip non-alphanumeric (keep hyphens + underscores)
    - max 64 chars

    Applied at the API boundary when a professor adds/renames a category so
    "Problem Sets" + "problem sets" + "Problem-Sets" all collapse to
    `problem_sets`. Prevents schema explosion from minor typing variance.
    """
    if not value:
        return ""
    s = value.strip().lower()
    s = re.sub(r"\s+", "_", s)
    s = re.sub(r"[^a-z0-9_\-]", "", s)
    return s[:64]


def classify_doc_category(text: str, filename: str, ta_categories: list) -> tuple:
    """Classify a document into one of the TA's configurable categories.

    Args:
        text: full document text
        filename: original filename
        ta_categories: list of {slug, label} dicts from TeachingAssistant.doc_categories

    Returns (slug, confidence, rationale) where:
      - slug: one of the slugs from ta_categories (always valid — falls back to
        the "other"-style category if no good match)
      - confidence: float in [0, 1]
      - rationale: short string explaining the choice

    The classifier sees both slug AND label for each candidate (research
    2026-05-22, Q2 — labels are user-meaningful, slugs are stable). The LLM
    picks the SLUG, which is what we persist.

    The Stage 2B per-TA configurable classifier. (Replaced an earlier
    hard-coded-enum classifier — `classify_doc_role` — which was deleted
    in Phase 1 cleanup since it had zero readers in retrieval.) See
    attached_assets/maize-retrieval-doc-classification-research-2026-05-22.md.
    """
    from openai import OpenAI

    if not ta_categories:
        # No TA categories configured — degrade gracefully. Caller should seed
        # defaults at TA creation, but we don't want to crash if they didn't.
        return "other", 0.0, "TA has no doc_categories configured"

    # Pick a fallback slug — prefer "other" if it exists, else the last entry.
    fallback_slug = next(
        (c["slug"] for c in ta_categories if c["slug"] == "other"),
        ta_categories[-1]["slug"],
    )
    valid_slugs = {c["slug"] for c in ta_categories}

    client = OpenAI(api_key=Config.OPENAI_API_KEY)
    preview = text[:3000] if len(text) > 3000 else text

    category_lines = "\n".join(
        f'- slug: "{c["slug"]}" — label: "{c["label"]}"' for c in ta_categories
    )

    prompt = f"""You are classifying a course document into the professor's configured category list for this teaching assistant. Each category has a stable internal SLUG and a human-readable LABEL. You will return the SLUG.

Filename: {filename}

Document preview:
{preview}

Available categories for this course:
{category_lines}

Pick the ONE category whose label best describes this document's purpose. Use the label to understand each category's meaning (the slug is just a stable identifier).

Examples of correct routing:
- A document containing problems for students to solve, named like "Quiz 2.pdf" → pick the category whose label matches "Quizzes" (or the closest match: "Homeworks", "Problem Sets", "Labs"). NOT solutions, even if it contains some worked examples — primary purpose matters.
- A solutions key like "Quiz 2 solutions.pdf" → pick the "Solutions" category if it exists, else the closest match.
- Interactive lecture slides → pick "Lectures".
- A course syllabus → pick "Syllabus".
- A formula sheet → pick "Reference Materials".
- If genuinely nothing fits → pick the catch-all (often "Other").

Return ONLY JSON in this exact shape:
{{"slug": "quizzes", "confidence": 0.92, "rationale": "Contains five problem statements consistent with a quiz; no worked answers shown."}}

Confidence: 0.9+ = certain; 0.7-0.9 = likely; <0.7 = uncertain. The slug MUST match one from the available list above exactly."""

    for attempt in range(3):
        try:
            response = client.chat.completions.create(
                model=Config.VISION_MODEL,
                messages=[{"role": "user", "content": prompt}],
                max_tokens=200,
                response_format={"type": "json_object"},
            )
            content = response.choices[0].message.content.strip()
            data = json.loads(content)
            slug = (data.get("slug") or fallback_slug).strip().lower()
            if slug not in valid_slugs:
                logger.warning(
                    f"classify_doc_category: LLM returned slug {slug!r} not in TA's list; "
                    f"falling back to {fallback_slug!r}"
                )
                slug = fallback_slug
            confidence = float(data.get("confidence") or 0.0)
            confidence = max(0.0, min(1.0, confidence))
            rationale = str(data.get("rationale") or "")[:500]
            return slug, confidence, rationale
        except Exception as e:
            logger.warning(f"classify_doc_category attempt {attempt+1} failed: {e}")
            if attempt == 2:
                logger.error("All classify_doc_category attempts failed; defaulting to fallback slug")
                return fallback_slug, 0.0, f"classification failed: {type(e).__name__}"


def summarize_doc(text: str, filename: str, content_title: str = "") -> str:
    """Phase B Stage B10: per-doc summary used by the future hybrid_doc_search refactor.

    Returns a ~100-150 word summary describing what the doc is, what it covers,
    where it sits in the course (week N / lesson topic / etc.), and key concepts.
    Caller embeds the returned text and persists to `Document.summary` +
    `Document.summary_embedding`.

    Indexing-only TODAY: nothing reads `Document.summary_embedding` from
    retrieval yet. The future Phase-B refactor (replacing BM25+dense+filename
    RRF in hybrid_doc_search with summary-cosine + LLM tiebreaker) is gated on
    these columns being populated first. See
    `attached_assets/maize-architecture-review-2026-05-23.md` section 3.3.

    Returns an empty string if the LLM call fails after 3 retries. Callers
    should treat empty as "no summary available" and leave the column NULL.
    """
    from openai import OpenAI

    if not (text and text.strip()):
        return ""

    client = OpenAI(api_key=Config.OPENAI_API_KEY)
    preview = text[:8000] if len(text) > 8000 else text  # generous window for richer summaries

    title_hint = f"\nContent title (from metadata): {content_title}" if content_title else ""

    prompt = f"""Summarize this course document in 100-150 words for a teaching-assistant retrieval system. The summary becomes the document's representation in a doc-routing index — when a student asks a question, the system will compare the question's embedding to your summary's embedding to decide whether this doc is relevant.

Filename: {filename}{title_hint}

Document text (truncated to first 8000 chars):
{preview}

Your summary must:
- State what KIND of document this is (lecture notes, problem set, quiz, syllabus, solutions, exam, reference sheet, etc.)
- State the doc's POSITION in the course if extractable (Week N, Lesson N, Lecture N, Quiz N, Pset N, course-level admin, etc.)
- List the 3-6 KEY TOPICS or CONCEPTS the doc covers
- Note any structural features that matter for routing (e.g. "covers multiple sections each with their own sub-problems", "single problem with sub-parts a-d", "slide deck with N slides on topic X")

Write a single paragraph. No bullet points. No markdown. No preamble like "This document is...". Start with the document's KIND and POSITION, then topics. ~100-150 words.

Return ONLY the summary text — no JSON wrapping, no surrounding quotes."""

    for attempt in range(3):
        try:
            response = client.chat.completions.create(
                model=Config.VISION_MODEL,  # gpt-4o; cheap + fast for summarization
                messages=[{"role": "user", "content": prompt}],
                max_tokens=300,
            )
            summary = (response.choices[0].message.content or "").strip()
            # Strip surrounding quotes if the LLM added them despite instructions
            if len(summary) >= 2 and summary[0] in ('"', "'") and summary[-1] == summary[0]:
                summary = summary[1:-1].strip()
            return summary
        except Exception as e:
            logger.warning(f"summarize_doc attempt {attempt+1} failed: {e}")
            if attempt == 2:
                logger.error(f"All summarize_doc attempts failed for {filename!r}; returning empty summary")
                return ""


def infer_doc_metadata_from_filename(filename: str) -> dict:
    """
    Synchronous lightweight metadata inference from a filename. Used at upload
    time so the docs table immediately shows a sensible doc_type instead of
    blank — the slow LLM-based extractor still runs in the background to refine.

    Returns a dict with whichever of {doc_type, assignment_number,
    instructional_unit_number} could be inferred. Empty dict if nothing matched.

    Patterns mirror analyze_query() in src/retriever.py for consistency between
    upload-time tagging and query-time filtering.
    """
    if not filename:
        return {}

    name_lower = filename.lower()
    # Strip extension for cleaner matching
    if '.' in name_lower:
        name_lower = name_lower.rsplit('.', 1)[0]

    out = {}

    # Capture 1-3 digit numbers only — avoids matching year-like values (2024, etc.)
    # as assignment/unit numbers.
    NUM = r'(\d{1,3})(?!\d)'

    # Homework / problem set — highest specificity first
    hw_match = re.search(r'(?:homework|hw|assignment|problem\s*set|p\s*set|pset|ps)\s*[-_#]?\s*' + NUM, name_lower)
    if hw_match:
        out['doc_type'] = 'homework'
        out['assignment_number'] = hw_match.group(1)
        return out
    if re.search(r'(?:^|[^a-z])(?:homework|hw|assignment|problem\s*set|pset)(?:[^a-z]|$)', name_lower):
        out['doc_type'] = 'homework'
        return out

    # Solutions to homework (still classify as homework — solutions ARE homework material)
    if re.search(r'(?:solution|answer|key)', name_lower) and re.search(r'(?:homework|hw|problem\s*set|pset|ps|assignment)', name_lower):
        out['doc_type'] = 'homework'
        sol_num = re.search(r'(?:homework|hw|problem\s*set|pset|ps|assignment)\s*[-_#]?\s*' + NUM, name_lower)
        if sol_num:
            out['assignment_number'] = sol_num.group(1)
        return out

    # Exam / quiz / midterm / final
    exam_num_match = re.search(r'(?:exam|midterm|quiz|test)\s*[-_#]?\s*' + NUM, name_lower)
    if exam_num_match:
        out['doc_type'] = 'exam'
        out['assignment_number'] = exam_num_match.group(1)
        return out
    if re.search(r'(?:^|[^a-z])(?:final|midterm|exam|quiz)(?:[^a-z]|$)', name_lower):
        out['doc_type'] = 'exam'
        return out

    # Lecture / class / week / module / session
    lec_match = re.search(r'(?:lecture|lec|class|week|module|session|day)\s*[-_#]?\s*' + NUM, name_lower)
    if lec_match:
        out['doc_type'] = 'lecture'
        out['instructional_unit_number'] = int(lec_match.group(1))
        return out
    if re.search(r'(?:^|[^a-z])(?:lecture|slides|notes)(?:[^a-z]|$)', name_lower):
        out['doc_type'] = 'lecture'
        return out

    # Reading / chapter / article
    ch_match = re.search(r'(?:chapter|ch|reading)\s*[-_#]?\s*' + NUM, name_lower)
    if ch_match:
        out['doc_type'] = 'reading'
        out['instructional_unit_number'] = int(ch_match.group(1))
        return out
    if re.search(r'(?:^|[^a-z])(?:reading|article|chapter|paper)(?:[^a-z]|$)', name_lower):
        out['doc_type'] = 'reading'
        return out

    # Syllabus
    if re.search(r'syllabus|course\s*outline', name_lower):
        out['doc_type'] = 'syllabus'
        return out

    return out


def extract_metadata_from_file_content(file_content: bytes, file_type: str, original_filename: str) -> dict:
    """
    Extract document metadata from file content at upload time.
    This runs LLM classification immediately so admins can review/edit before indexing.

    Returns dict with: doc_type, assignment_number, instructional_unit_number, content_title, etc.
    """
    import tempfile
    
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_type}") as tmp_file:
            tmp_file.write(file_content)
            tmp_path = tmp_file.name
        
        try:
            text, _ = extract_text_from_file(tmp_path)
        finally:
            os.unlink(tmp_path)
        
        if not text or len(text.strip()) < 50:
            logger.warning(f"Could not extract sufficient text from {original_filename} for metadata extraction")
            return {
                "doc_type": None,
                "assignment_number": None,
                "instructional_unit_number": None,
                "instructional_unit_label": None,
                "content_title": None,
                "extraction_success": False
            }
        
        metadata = extract_metadata_with_llm(text, original_filename)
        metadata["extraction_success"] = True
        return metadata
        
    except Exception as e:
        logger.error(f"Error extracting metadata from {original_filename}: {e}")
        return {
            "doc_type": None,
            "assignment_number": None,
            "instructional_unit_number": None,
            "instructional_unit_label": None,
            "content_title": None,
            "extraction_success": False
        }


def chunk_text(text: str, chunk_size: int = 800, overlap: int = 200) -> list:
    if len(text) <= chunk_size:
        return [text]
    
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        
        if end < len(text):
            break_points = [
                text.rfind('\n\n', start, end),
                text.rfind('. ', start, end),
                text.rfind('\n', start, end),
                text.rfind(' ', start, end)
            ]
            for bp in break_points:
                if bp > start + chunk_size // 2:
                    end = bp + 1
                    break
        
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        
        start = end - overlap
        if start < 0:
            start = 0
        if start >= len(text):
            break
    
    return chunks

def process_and_index_documents_resumable(ta_id: str, progress_callback=None, resume_from_doc_id=None, job_start_time=None) -> dict:
    """
    Resumable version of document indexing that can continue from where it left off.
    Uses document.last_indexed_at as the completion marker - if a document has chunks
    stored, it has last_indexed_at set and will be skipped on resume.
    """
    import tempfile
    from openai import OpenAI
    from datetime import datetime
    
    from models import db, Document, TeachingAssistant, DocumentChunk
    from flask import current_app
    from src.qa_logger import log_index_batch
    
    is_resume = bool(resume_from_doc_id)
    logger.info(f"[{ta_id}] Starting {'resumable' if is_resume else 'fresh'} indexing process...")
    
    ta = db.session.get(TeachingAssistant, ta_id)
    ta_slug = ta.slug if ta else "unknown"
    
    client = OpenAI(api_key=Config.OPENAI_API_KEY)
    
    all_doc_ids = [d.id for d in db.session.query(Document.id).filter_by(ta_id=ta_id).order_by(Document.id).all()]
    total_docs = len(all_doc_ids)
    
    if total_docs == 0:
        raise ValueError("No documents found for this TA")
    
    if is_resume:
        # Skip docs that already have chunks. `last_indexed_at IS NULL` is the
        # trigger for "needs processing" — set when the doc is freshly uploaded
        # and cleared on full-rebuild paths. (We deliberately don't compare
        # against `updated_at` because metadata edits like display_name or
        # doc_type bump it without invalidating the chunk content.)
        already_indexed_doc_ids = set(
            d.id for d in db.session.query(Document.id).filter(
                Document.ta_id == ta_id,
                Document.last_indexed_at.isnot(None),
            ).all()
        )
        doc_ids = [d for d in all_doc_ids if d not in already_indexed_doc_ids]
        docs_already_processed = len(already_indexed_doc_ids)
        logger.info(f"[{ta_id}] Incremental indexing: {docs_already_processed} docs already indexed, {len(doc_ids)} remaining")
    else:
        logger.info(f"[{ta_id}] Fresh indexing - clearing existing chunks and reset last_indexed_at...")
        DocumentChunk.query.filter_by(ta_id=ta_id).delete()
        Document.query.filter_by(ta_id=ta_id).update({"last_indexed_at": None}, synchronize_session=False)
        db_commit_with_retry(db)
        logger.info(f"[{ta_id}] Cleared existing chunks and reset document states")
        doc_ids = all_doc_ids
        docs_already_processed = 0
    
    logger.info(f"[{ta_id}] Found {len(doc_ids)} documents to process: {doc_ids}")
    
    total_chunks_created = 0
    all_index_log_entries = []
    docs_with_content = 0
    docs_succeeded: list[dict] = []
    docs_failed: list[dict] = []

    for doc_idx, doc_id in enumerate(doc_ids):
        absolute_doc_idx = docs_already_processed + doc_idx

        doc = db.session.get(Document, doc_id)
        if not doc:
            logger.warning(f"[{ta_id}] Document {doc_id} not found, skipping")
            docs_failed.append({
                "doc_id": doc_id,
                "filename": "<not found>",
                "error": "Document row missing at index time — may have been deleted mid-run.",
            })
            continue

        logger.info(f"[{ta_id}] Processing document [{doc.id}]: {doc.original_filename} ({absolute_doc_idx + 1}/{total_docs})")

        if progress_callback and total_docs > 0:
            progress = int((absolute_doc_idx / total_docs) * 80)
            progress_callback(ta_id, progress, docs_processed=absolute_doc_idx)

        text = None
        page_count = 0

        # Keep the indexing watchdog informed WHILE a single document is being
        # extracted. check_stale_indexing_jobs fails any job whose row has not been
        # touched in 5 minutes, and a long PDF's page-by-page vision work can run
        # well past that inside one document — the job is progressing, just
        # silently. Re-reports the current document's progress, which is enough to
        # move IndexingJob.updated_at (see update_indexing_progress in app.py).
        _hb_progress = int((absolute_doc_idx / total_docs) * 80) if total_docs else 0

        def _heartbeat(_p=_hb_progress, _i=absolute_doc_idx):
            if not progress_callback:
                return
            try:
                progress_callback(ta_id, _p, docs_processed=_i)
            except Exception as hb_e:
                logger.warning(f"[{ta_id}] heartbeat failed: {type(hb_e).__name__}: {hb_e}")

        logger.info(f"[{ta_id}] [{doc.id}] Extracting text...")
        if doc.file_content:
            with tempfile.NamedTemporaryFile(delete=False, suffix=f".{doc.file_type}") as tmp_file:
                tmp_file.write(doc.file_content)
                tmp_path = tmp_file.name
            try:
                text, page_count = extract_text_from_file(tmp_path, heartbeat=_heartbeat)
            finally:
                os.unlink(tmp_path)
        elif os.path.exists(doc.storage_path):
            text, page_count = extract_text_from_file(doc.storage_path, heartbeat=_heartbeat)
        else:
            err = "No file content available — please re-upload this document."
            logger.warning(f"[{ta_id}] [{doc.id}] {err}")
            doc.extraction_metadata = {
                "_indexing_status": "extraction_failed",
                "_error": err,
            }
            db_commit_with_retry(db)
            docs_failed.append({"doc_id": doc.id, "filename": doc.original_filename, "error": err})
            continue

        if not text:
            err = (
                f"Text extraction returned empty for .{doc.file_type or 'unknown'} file. "
                f"The file may be corrupt, image-only, password-protected, or in an unsupported variant."
            )
            logger.warning(f"[{ta_id}] [{doc.id}] {err}")
            doc.extraction_metadata = {
                "_indexing_status": "extraction_failed",
                "_error": err,
            }
            db_commit_with_retry(db)
            docs_failed.append({"doc_id": doc.id, "filename": doc.original_filename, "error": err})
            continue
        
        raw_text_length = len(text)
        logger.info(f"[{ta_id}] [{doc.id}] Extracted {raw_text_length} chars from {page_count} pages")
        
        logger.info(f"[{ta_id}] [{doc.id}] Extracting metadata with LLM...")
        metadata = extract_metadata_with_llm(text, doc.original_filename)
        
        if not doc.doc_type:
            doc.doc_type = metadata.get("doc_type")
        if not doc.assignment_number:
            doc.assignment_number = metadata.get("assignment_number")
        if doc.instructional_unit_number is None:
            doc.instructional_unit_number = metadata.get("instructional_unit_number")
        if not doc.instructional_unit_label:
            doc.instructional_unit_label = metadata.get("instructional_unit_label")
        if not doc.content_title:
            doc.content_title = metadata.get("content_title")
        
        doc.extraction_metadata = metadata
        doc.metadata_extracted = True

        # Phase A Stage 2B (research 2026-05-22): classify doc_category using
        # the parent TA's per-tenant configurable list. PRIMARY axis for
        # retrieval. Skip if already set (the UI/PATCH route is the only
        # other writer and we treat its value as authoritative).
        ta_categories = (ta.doc_categories if ta else None) or []
        if not doc.doc_category and ta_categories:
            logger.info(f"[{ta_id}] [{doc.id}] Classifying doc_category against {len(ta_categories)} TA categories...")
            try:
                slug, cat_conf, cat_rationale = classify_doc_category(
                    text, doc.original_filename, ta_categories
                )
                doc.doc_category = slug
                logger.info(f"[{ta_id}] [{doc.id}] doc_category={slug} (conf={cat_conf:.2f})")
            except Exception as cat_err:
                logger.warning(f"[{ta_id}] [{doc.id}] doc_category classification failed: {cat_err}; leaving unset")
        elif not ta_categories:
            logger.warning(f"[{ta_id}] [{doc.id}] TA has no doc_categories — skipping classify_doc_category")

        # Phase B Stage B10: per-doc LLM-generated summary + summary embedding.
        # Indexing-only today (no retrieval reads summary_embedding yet); sets up
        # the future hybrid_doc_search refactor. Skips when summary is already
        # populated AND not flagged stale (no stale flag today, so set-once
        # semantics: rebuild a stale summary by manually clearing the column
        # OR running backfill --force). See attached_assets/maize-architecture-review-2026-05-23.md.
        if not doc.summary:
            logger.info(f"[{ta_id}] [{doc.id}] Generating doc summary...")
            try:
                summary_text = summarize_doc(text, doc.original_filename, doc.content_title or "")
                if summary_text:
                    summary_emb_response = client.embeddings.create(
                        model=Config.EMBEDDING_MODEL,
                        input=summary_text,
                    )
                    doc.summary = summary_text
                    doc.summary_embedding = summary_emb_response.data[0].embedding
                    logger.info(f"[{ta_id}] [{doc.id}] Summary populated ({len(summary_text)} chars, ~{len(summary_text.split())} words)")
                else:
                    logger.warning(f"[{ta_id}] [{doc.id}] summarize_doc returned empty; leaving summary unset")
            except Exception as summary_err:
                logger.warning(f"[{ta_id}] [{doc.id}] Summary generation failed: {summary_err}; leaving unset")

        # Phase B latency Phase 1 (2026-08-06). Cache the extracted full text
        # so the hybrid_full_doc fallback (src/retriever.py:get_full_document_text)
        # can serve it without re-running pdfplumber + gpt-4o vision extraction.
        # Pilot data (July 2026) showed re-extraction costing 30-40s per fallback
        # fire — this caches the result once at indexing time. Set-once: skips
        # when already populated (backfill --force to rebuild).
        if not doc.full_text:
            doc.full_text = sanitize_text(text)
            logger.info(f"[{ta_id}] [{doc.id}] Cached full_text ({len(doc.full_text)} chars)")

        # Phase A retrieval refactor: BM25 tsvector for hybrid Stage 1 retrieval.
        # Built from the already-extracted text via PostgreSQL's to_tsvector.
        doc.bm25_tsvector = db.func.to_tsvector('english', sanitize_text(text))

        logger.info(f"[{ta_id}] [{doc.id}] Saving metadata (preserving human edits)...")
        db_commit_with_retry(db)
        logger.info(f"[{ta_id}] [{doc.id}] Metadata saved")
        
        headers_found = extract_section_headers(text)
        headers_summary = "; ".join([f"{h[1][:40]}" for h in headers_found[:5]])
        logger.info(f"[{ta_id}] [{doc.id}] Found {len(headers_found)} headers: {headers_summary}")
        
        chunks = chunk_text_with_context(text, Config.CHUNK_SIZE, Config.CHUNK_OVERLAP, doc.original_filename)
        num_chunks = len(chunks)
        
        doc_chunk_data = []
        doc_log_entries = []
        for i, chunk_data in enumerate(chunks):
            doc_chunk_data.append({
                "chunk_index": i,
                "chunk_text": chunk_data["original_text"],
                "chunk_text_enriched": chunk_data["text"],
                "context": chunk_data.get("context", ""),
                "section_path": chunk_data.get("section_path") or [],  # Phase B B8
                "doc_type": doc.doc_type or "other",
                "assignment_number": doc.assignment_number or "",
                "instructional_unit_number": doc.instructional_unit_number or 0,
                "instructional_unit_label": doc.instructional_unit_label or "",
                "file_name": doc.display_name or doc.original_filename,
                "doc_role": doc.doc_role,  # Phase A — populated above; may be None on legacy paths
                "doc_category": doc.doc_category,  # Phase A Stage 2B — populated above; may be None on legacy paths
            })
            
            doc_log_entries.append({
                "ta_id": ta_id,
                "ta_slug": ta_slug,
                "file_name": doc.display_name or doc.original_filename,
                "doc_type": doc.doc_type or "other",
                "total_pages": page_count,
                "raw_text_length": raw_text_length,
                "chunk_index": i,
                "total_chunks": num_chunks,
                "chunk_text_length": len(chunk_data["original_text"]),
                "chunk_context": chunk_data.get("context", ""),
                "chunk_text_preview": chunk_data["original_text"][:300],
                "enriched_text_preview": chunk_data["text"][:300],
                "has_embedding": False,
                "status": "pending",
                "error_message": "",
                "headers_found": headers_summary if i == 0 else ""
            })
        
        try:
            logger.info(f"[{ta_id}] [{doc.id}] Embedding {num_chunks} chunks for this document...")
            chunk_texts = [c["chunk_text_enriched"] for c in doc_chunk_data]
            
            doc_embeddings = []
            batch_size = 100
            for batch_start in range(0, len(chunk_texts), batch_size):
                batch_texts_slice = chunk_texts[batch_start:batch_start+batch_size]
                response = client.embeddings.create(
                    model=Config.EMBEDDING_MODEL,
                    input=batch_texts_slice
                )
                doc_embeddings.extend([item.embedding for item in response.data])
            
            logger.info(f"[{ta_id}] [{doc.id}] Storing {num_chunks} chunks...")
            for i, chunk_item in enumerate(doc_chunk_data):
                chunk_obj = DocumentChunk(
                    ta_id=ta_id,
                    document_id=doc.id,
                    chunk_index=chunk_item["chunk_index"],
                    chunk_text=sanitize_text(chunk_item["chunk_text"]),
                    chunk_context=chunk_item.get("context", "")[:256] if chunk_item.get("context") else None,
                    section_path=chunk_item.get("section_path") or [],  # Phase B B8
                    doc_type=chunk_item["doc_type"],
                    assignment_number=chunk_item["assignment_number"],
                    instructional_unit_number=chunk_item["instructional_unit_number"],
                    instructional_unit_label=chunk_item["instructional_unit_label"],
                    file_name=chunk_item["file_name"],
                    doc_role=chunk_item.get("doc_role"),  # Phase A — None on legacy paths is OK
                    doc_category=chunk_item.get("doc_category"),  # Phase A Stage 2B
                    embedding=doc_embeddings[i]
                )
                db.session.add(chunk_obj)
            
            now = datetime.utcnow()
            doc.last_indexed_at = now
            doc.updated_at = now
            db_commit_with_retry(db)

            total_chunks_created += num_chunks
            docs_with_content += 1
            docs_succeeded.append({
                "doc_id": doc.id,
                "filename": doc.original_filename,
                "chunks": num_chunks,
            })
            logger.info(f"[{ta_id}] [{doc.id}] Document fully indexed with {num_chunks} chunks")

            for entry in doc_log_entries:
                entry["has_embedding"] = True
                entry["status"] = "success"
            all_index_log_entries.extend(doc_log_entries)

        except Exception as doc_error:
            logger.error(f"[{ta_id}] [{doc.id}] Failed to embed/store chunks: {doc_error}")
            db.session.rollback()
            for entry in doc_log_entries:
                entry["has_embedding"] = False
                entry["status"] = "error"
                entry["error_message"] = str(doc_error)[:500]
            all_index_log_entries.extend(doc_log_entries)
            docs_failed.append({
                "doc_id": doc.id,
                "filename": doc.original_filename,
                "error": f"Embedding/storage failed: {str(doc_error)[:300]}",
            })
            raise
        
        if progress_callback:
            progress_callback(ta_id, int(((absolute_doc_idx + 1) / total_docs) * 100), 
                            docs_processed=absolute_doc_idx + 1,
                            chunks_created=total_chunks_created)
    
    if docs_with_content == 0:
        if is_resume:
            logger.info(f"[{ta_id}] No new documents to process - resumption complete")
            return {
                "chunks_indexed": 0,
                "docs_succeeded": docs_succeeded,
                "docs_failed": docs_failed,
            }
        if not docs_failed:
            raise ValueError("No text content found in any documents")
        # All processed docs failed extraction; surface per-doc errors but don't crash the run.
        logger.warning(f"[{ta_id}] All {len(docs_failed)} processed docs failed extraction")

    logger.info(
        f"[{ta_id}] Indexing complete! Total chunks: {total_chunks_created}, "
        f"succeeded: {len(docs_succeeded)}, failed: {len(docs_failed)}"
    )

    if all_index_log_entries:
        logger.info(f"[{ta_id}] Logging {len(all_index_log_entries)} index entries to Google Sheets...")
        log_index_batch(all_index_log_entries)

    return {
        "chunks_indexed": total_chunks_created,
        "docs_succeeded": docs_succeeded,
        "docs_failed": docs_failed,
    }
